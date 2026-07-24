"""Professional templates for LibreOffice MCP.

Each template is a function that takes user data and generates a polished
file using the existing presentation/document/spreadsheet primitives.
"""
from __future__ import annotations

import json
import os
import tempfile
from pathlib import Path
from typing import Any

from libreoffice_mcp.presentation import (
    create_presentation,
    add_slide,
    set_slide_title,
    add_bullet_points,
    add_table,
)
from libreoffice_mcp.document import (
    create_document,
    add_heading,
    add_paragraph,
)
from libreoffice_mcp.spreadsheet import (
    create_spreadsheet,
    add_sheet,
    set_cell_value,
)


def executive_summary_presentation(output_path: str, data: dict[str, Any]) -> None:
    """Create a corporate executive summary presentation.
    
    Expected data keys:
      - title: str (main presentation title)
      - subtitle: str
      - presenter: str
      - date: str
      - metrics: list of {label, value, description}
      - financials: {headers, rows} for table
      - roadmap: list of {phase, quarter, details}
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    
    DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
    ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
    LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    def bg(slide, shape_type, left, top, width, height, fill_color):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        spTree = slide.shapes._spTree
        sp = shape._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    
    # Slide 1: Hero Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(3), DARK_BLUE)
    bg(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(3), Inches(13.333), Inches(0.05), ACCENT_ORANGE)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = data.get("title", "Executive Summary")
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1))
    tf2 = sub_box.text_frame
    tf2.text = data.get("subtitle", "")
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.size = Pt(24)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    
    info = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1))
    tf3 = info.text_frame
    tf3.text = f"{data.get('presenter', '')}\n{data.get('date', '')}"
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.runs[0]
    run3.font.size = Pt(16)
    run3.font.color.rgb = DARK_TEXT
    
    # Slide 2: Key Metrics (3-column)
    if data.get("metrics"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), DARK_BLUE)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "Key Metrics"
        p = tf.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        metrics = data["metrics"][:3]  # Max 3 for clean layout
        for i, m in enumerate(metrics):
            left = Inches(0.5 + i * 4.2)
            bg(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(1.2), LIGHT_GRAY)
            
            tb = slide.shapes.add_textbox(left, Inches(1.6), Inches(3.8), Inches(0.4))
            tf = tb.text_frame
            tf.text = m.get("label", "")
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = DARK_BLUE
            
            mb = slide.shapes.add_textbox(left, Inches(2.0), Inches(3.8), Inches(0.5))
            tf2 = mb.text_frame
            tf2.text = m.get("value", "")
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.runs[0]
            run2.font.size = Pt(28)
            run2.font.bold = True
            run2.font.color.rgb = ACCENT_ORANGE
            
            db = slide.shapes.add_textbox(left + Inches(0.1), Inches(3.0), Inches(3.6), Inches(1.5))
            tf3 = db.text_frame
            tf3.word_wrap = True
            tf3.text = m.get("description", "")
            p3 = tf3.paragraphs[0]
            p3.alignment = PP_ALIGN.CENTER
            run3 = p3.runs[0]
            run3.font.size = Pt(14)
            run3.font.color.rgb = DARK_TEXT
    
    # Slide 3: Financial Table
    if data.get("financials"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), DARK_BLUE)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "Financial Projections"
        p = tf.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        fin = data["financials"]
        headers = fin.get("headers", [])
        rows = fin.get("rows", [])
        num_rows = 1 + len(rows)
        num_cols = len(headers)
        
        table = slide.shapes.add_table(num_rows, num_cols, Inches(1.5), Inches(1.5), Inches(10), Inches(4)).table
        col_width = Inches(10 / num_cols)
        for col in table.columns:
            col.width = col_width
        
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(16)
                run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
        
        for i, row_data in enumerate(rows):
            for j, val in enumerate(row_data):
                cell = table.cell(i + 1, j)
                cell.text = val
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(14)
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # Slide 4: Roadmap Timeline
    if data.get("roadmap"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8), DARK_BLUE)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "Strategic Roadmap"
        p = tf.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(11.333), Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        bar.line.fill.background()
        
        phases = data["roadmap"][:3]
        for i, phase in enumerate(phases):
            left = Inches(1.5 + i * 3.8)
            
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.5), Inches(2.35), Inches(0.4), Inches(0.4))
            dot.fill.solid()
            dot.fill.fore_color.rgb = ACCENT_ORANGE if i == 1 else DARK_BLUE
            dot.line.color.rgb = WHITE
            dot.line.width = Pt(2)
            
            tb = slide.shapes.add_textbox(left, Inches(1.2), Inches(3.5), Inches(0.8))
            tf = tb.text_frame
            tf.text = phase.get("phase", "")
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = DARK_BLUE
            
            qb = slide.shapes.add_textbox(left, Inches(3.2), Inches(3.5), Inches(0.4))
            tf2 = qb.text_frame
            tf2.text = phase.get("quarter", "")
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.runs[0]
            run2.font.size = Pt(14)
            run2.font.bold = True
            run2.font.color.rgb = ACCENT_ORANGE
            
            db = slide.shapes.add_textbox(left, Inches(3.8), Inches(3.5), Inches(2.5))
            tf3 = db.text_frame
            tf3.word_wrap = True
            tf3.text = phase.get("details", "")
            for p3 in tf3.paragraphs:
                p3.alignment = PP_ALIGN.CENTER
                for run3 in p3.runs:
                    run3.font.size = Pt(13)
                    run3.font.color.rgb = DARK_TEXT
    
    # Slide 5: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BLUE)
    
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.166), Inches(1.5), Inches(3), Inches(3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_ORANGE
    circle.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.333), Inches(1.5))
    tf = tb.text_frame
    tf.text = "Thank You"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(66)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    st = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    stf = st.text_frame
    stf.text = "Questions & Discussion"
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.runs[0]
    sr.font.size = Pt(28)
    sr.font.italic = True
    sr.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    
    ct = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.8))
    ctf = ct.text_frame
    ctf.text = data.get("presenter", "")
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.runs[0]
    cr.font.size = Pt(14)
    cr.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    
    prs.save(output_path)


def project_proposal_document(output_path: str, data: dict[str, Any]) -> None:
    """Create a professional project proposal document.
    
    Expected data keys:
      - title: str
      - client: str
      - date: str
      - author: str
      - sections: list of {heading, paragraphs}
      - budget_table: {headers, rows}
    """
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    doc = Document()
    
    # Title page
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_para.add_run(data.get("title", "Project Proposal"))
    run.font.size = Pt(28)
    run.font.bold = True
    
    doc.add_paragraph()
    
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"Prepared for: {data.get('client', '')}\n").font.size = Pt(14)
    meta.add_run(f"Prepared by: {data.get('author', '')}\n").font.size = Pt(14)
    meta.add_run(f"Date: {data.get('date', '')}\n").font.size = Pt(14)
    
    doc.add_page_break()
    
    # Executive Summary
    doc.add_heading("Executive Summary", level=1)
    if data.get("summary"):
        doc.add_paragraph(data["summary"])
    
    # Sections
    for section in data.get("sections", []):
        doc.add_heading(section.get("heading", ""), level=2)
        for para_text in section.get("paragraphs", []):
            doc.add_paragraph(para_text)
    
    # Budget table
    if data.get("budget_table"):
        doc.add_heading("Budget Overview", level=2)
        table = doc.add_table(rows=1, cols=len(data["budget_table"]["headers"]))
        table.style = "Table Grid"
        
        hdr_cells = table.rows[0].cells
        for j, header in enumerate(data["budget_table"]["headers"]):
            hdr_cells[j].text = header
            for paragraph in hdr_cells[j].paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
        
        for row_data in data["budget_table"]["rows"]:
            row_cells = table.add_row().cells
            for j, val in enumerate(row_data):
                row_cells[j].text = val
    
    doc.save(output_path)


def financial_report_spreadsheet(output_path: str, data: dict[str, Any]) -> None:
    """Create a financial report spreadsheet with charts.
    
    Expected data keys:
      - title: str
      - months: list of str (month labels)
      - revenue: list of numbers
      - expenses: list of numbers
      - profit: list of numbers
    """
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.styles import Font, PatternFill, Alignment
    
    wb = Workbook()
    ws = wb.active
    ws.title = "Financial Report"
    
    # Header
    ws["A1"] = data.get("title", "Financial Report")
    ws["A1"].font = Font(size=18, bold=True)
    ws.merge_cells("A1:E1")
    
    # Column headers
    headers = ["Month", "Revenue", "Expenses", "Profit", "Margin %"]
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=3, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="1A237E", end_color="1A237E", fill_type="solid")
        cell.alignment = Alignment(horizontal="center")
    
    # Data
    months = data.get("months", [])
    revenue = data.get("revenue", [])
    expenses = data.get("expenses", [])
    profit = data.get("profit", [])
    
    for i, month in enumerate(months):
        row = 4 + i
        ws.cell(row=row, column=1, value=month)
        ws.cell(row=row, column=2, value=revenue[i] if i < len(revenue) else 0)
        ws.cell(row=row, column=3, value=expenses[i] if i < len(expenses) else 0)
        ws.cell(row=row, column=4, value=profit[i] if i < len(profit) else 0)
        # Margin formula
        ws.cell(row=row, column=5, value=f"=D{row}/B{row}")
        ws.cell(row=row, column=5).number_format = "0.0%"
    
    # Chart
    last_row = 3 + len(months)
    chart = BarChart()
    chart.type = "col"
    chart.title = data.get("title", "Financial Report")
    chart.y_axis.title = "Amount ($)"
    chart.x_axis.title = "Month"
    
    data_ref = Reference(ws, min_col=2, min_row=3, max_col=4, max_row=last_row)
    cats_ref = Reference(ws, min_col=1, min_row=4, max_row=last_row)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)
    
    ws.add_chart(chart, "G3")
    
    wb.save(output_path)


TEMPLATES = {
    "executive_summary_presentation": executive_summary_presentation,
    "project_proposal_document": project_proposal_document,
    "financial_report_spreadsheet": financial_report_spreadsheet,
}
