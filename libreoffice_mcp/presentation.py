"""Presentation (pptx) operations for LibreOffice MCP."""
from __future__ import annotations

import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _layout_map(name: str) -> int:
    """Map friendly layout names to python-pptx layout indices."""
    mapping = {
        "title": 0,
        "title_and_content": 1,
        "blank": 6,
        "title_only": 5,
    }
    return mapping.get(name, 1)


def create_presentation(output_path: str | Path, *, width: float = 13.333, height: float = 7.5) -> None:
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    # Add a title slide by default
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(output_path))


def add_slide(file_path: str | Path, layout: str = "title_and_content") -> int:
    prs = Presentation(str(file_path))
    slide = prs.slides.add_slide(prs.slide_layouts[_layout_map(layout)])
    prs.save(str(file_path))
    return prs.slides.index(slide)


def set_slide_title(
    file_path: str | Path,
    slide_index: int,
    title: str,
    *,
    font_size: int | None = None,
    bold: bool | None = None,
) -> None:
    prs = Presentation(str(file_path))
    slide = prs.slides[slide_index]
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.text = title
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
    prs.save(str(file_path))


def add_bullet_points(
    file_path: str | Path,
    slide_index: int,
    bullets: list[str],
    *,
    font_size: int | None = None,
) -> None:
    prs = Presentation(str(file_path))
    slide = prs.slides[slide_index]
    # Find the first body placeholder
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    if body is None:
        # No body placeholder — add a text box
        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf = body.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = 0
        if font_size:
            for run in p.runs:
                run.font.size = Pt(font_size)
    prs.save(str(file_path))


def add_table(
    file_path: str | Path,
    slide_index: int,
    headers: list[str],
    rows: list[list[str]],
    *,
    left: float = 1.0,
    top: float = 2.0,
    width: float = 10.0,
    height: float = 3.0,
) -> None:
    prs = Presentation(str(file_path))
    slide = prs.slides[slide_index]
    num_rows = 1 + len(rows)
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(left), Inches(top), Inches(width), Inches(height)).table

    # Set column widths
    col_width = Inches(width / num_cols)
    for col in table.columns:
        col.width = col_width

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(14)

    # Data rows
    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(12)
        if i % 2 == 1:
            for j in range(num_cols):
                cell = table.cell(i + 1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    prs.save(str(file_path))


def add_image(
    file_path: str | Path,
    slide_index: int,
    image_path: str | Path,
    *,
    left: float = 1.0,
    top: float = 1.0,
    width: float | None = None,
    height: float | None = None,
) -> None:
    prs = Presentation(str(file_path))
    slide = prs.slides[slide_index]
    kwargs: dict[str, Any] = {"left": Inches(left), "top": Inches(top)}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(image_path), **kwargs)
    prs.save(str(file_path))


def export_pdf(file_path: str | Path, output_dir: str | None = None) -> Path:
    file_path = Path(file_path).resolve()
    out_dir = Path(output_dir).resolve() if output_dir else file_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        "soffice",
        "--headless",
        "--convertTo",
        "pdf",
        "--outdir",
        str(out_dir),
        str(file_path),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise RuntimeError(f"soffice failed: {result.stderr}")

    pdf_name = file_path.stem + ".pdf"
    pdf_path = out_dir / pdf_name
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF was not created at {pdf_path}")
    return pdf_path


def get_presentation_info(file_path: str | Path) -> dict[str, Any]:
    prs = Presentation(str(file_path))
    slides_info = []
    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text
        slides_info.append({"index": i, "title": title})
    return {
        "file": str(file_path),
        "slide_count": len(prs.slides),
        "width_inches": prs.slide_width / 914400,
        "height_inches": prs.slide_height / 914400,
        "slides": slides_info,
    }
